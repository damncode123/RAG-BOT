import datetime
import tempfile
import json
import csv
import xml.etree.ElementTree as ET
import logging
import os
from io import StringIO
from db.models import save_file_metadata
from api.notifications import send_notification

# Load environment variables early
from dotenv import load_dotenv
load_dotenv()

# External API credentials
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = os.getenv("PINECONE_ENV")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Document parsing libraries
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None
    logging.warning("PyPDF2 not installed. PDF parsing will not work.")

try:
    from docx import Document
except ImportError:
    Document = None
    logging.warning("python-docx not installed. DOCX parsing will not work.")

try:
    import openpyxl
except ImportError:
    openpyxl = None
    logging.warning("openpyxl not installed. Excel parsing will not work.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logging.warning("python-pptx not installed. PowerPoint parsing will not work.")

# LlamaIndex and Pinecone imports
from llama_index.embeddings.google import GeminiEmbedding


def get_pinecone_index():
    from pinecone import Pinecone
    pc = Pinecone(api_key=PINECONE_API_KEY)
    return pc.Index(PINECONE_INDEX_NAME)


def parse_text_file(contents, encoding='utf-8'):
    try:
        return contents.decode(encoding)
    except UnicodeDecodeError:
        for enc in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                return contents.decode(enc)
            except UnicodeDecodeError:
                continue
        return contents.decode('utf-8', errors='ignore')

def parse_markdown(contents):
    return parse_text_file(contents)

def parse_csv(contents):
    try:
        text_content = parse_text_file(contents)
        csv_reader = csv.reader(StringIO(text_content))
        rows = list(csv_reader)
        text_parts = []
        for i, row in enumerate(rows):
            if i == 0:
                text_parts.append(f"Headers: {', '.join(row)}")
            else:
                text_parts.append(f"Row {i}: {', '.join(row)}")
        return "\n".join(text_parts)
    except Exception as e:
        logging.error(f"Error parsing CSV: {e}")
        return parse_text_file(contents)

def parse_json(contents):
    try:
        data = json.loads(contents.decode('utf-8'))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        logging.error(f"Error parsing JSON: {e}")
        return parse_text_file(contents)

def parse_xml(contents):
    try:
        root = ET.fromstring(contents.decode('utf-8'))
        return ET.tostring(root, encoding='unicode', method='xml')
    except Exception as e:
        logging.error(f"Error parsing XML: {e}")
        return parse_text_file(contents)

def parse_html(contents):
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(contents.decode('utf-8'), 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        return soup.get_text()
    except ImportError:
        logging.warning("BeautifulSoup not installed. HTML parsing will be basic.")
        return parse_text_file(contents)
    except Exception as e:
        logging.error(f"Error parsing HTML: {e}")
        return parse_text_file(contents)

def parse_pdf(contents):
    if PdfReader is None:
        raise ImportError("PyPDF2 is required for PDF parsing")
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
        tmp.write(contents)
        tmp.flush()
        tmp_path = tmp.name
    try:
        reader = PdfReader(tmp_path)
        text_parts = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(text_parts)
    finally:
        os.remove(tmp_path)

def parse_docx(contents):
    if Document is None:
        raise ImportError("python-docx is required for DOCX parsing")
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
        tmp.write(contents)
        tmp.flush()
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                text_parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n\n".join(text_parts)
    finally:
        os.remove(tmp_path)

def parse_excel(contents, filename):
    if openpyxl is None:
        raise ImportError("openpyxl is required for Excel parsing")
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
        tmp.write(contents)
        tmp.flush()
        tmp_path = tmp.name
    try:
        workbook = openpyxl.load_workbook(tmp_path, data_only=True)
        text_parts = []
        for sheet in workbook.worksheets:
            sheet_text = [f"Sheet: {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell else "" for cell in row]
                sheet_text.append(" | ".join(row_text))
            text_parts.append("\n".join(sheet_text))
        return "\n\n".join(text_parts)
    finally:
        os.remove(tmp_path)

def parse_powerpoint(contents, filename):
    if Presentation is None:
        raise ImportError("python-pptx is required for PowerPoint parsing")
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
        tmp.write(contents)
        tmp.flush()
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"Slide {i+1}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_parts.append("\n".join(slide_text))
        return "\n\n".join(text_parts)
    finally:
        os.remove(tmp_path)

def parse_code_file(contents, filename):
    return parse_text_file(contents)

def parse_config_file(contents, filename):
    return parse_text_file(contents)

def parse_file(contents, filename):
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == '.txt': return parse_text_file(contents)
        elif ext == '.md': return parse_markdown(contents)
        elif ext == '.csv': return parse_csv(contents)
        elif ext == '.json': return parse_json(contents)
        elif ext == '.xml': return parse_xml(contents)
        elif ext in ['.html', '.htm']: return parse_html(contents)
        elif ext == '.pdf': return parse_pdf(contents)
        elif ext == '.docx': return parse_docx(contents)
        elif ext in ['.xls', '.xlsx']: return parse_excel(contents, filename)
        elif ext in ['.ppt', '.pptx']: return parse_powerpoint(contents, filename)
        elif ext in ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala', '.sql', '.sh', '.bat', '.ps1']:
            return parse_code_file(contents, filename)
        elif ext in ['.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf']:
            return parse_config_file(contents, filename)
        else:
            logging.warning(f"Unknown file type {ext}, treating as text")
            return parse_text_file(contents)
    except Exception as e:
        logging.error(f"Error parsing file {filename}: {e}")
        return parse_text_file(contents)

def chunk_text(text, chunk_size=300, overlap=50):
    if not text.strip(): return []
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size - overlap)]

def store_embeddings(user_id, filename, chunks):
    logging.info(f"Storing embeddings for {filename}, user {user_id}")
    embed_model = GeminiEmbedding(api_key=GEMINI_API_KEY, model_name="models/embedding-001")
    pinecone_index = get_pinecone_index()
    for i, chunk in enumerate(chunks):
        embedding = embed_model.get_text_embedding(chunk)
        pinecone_index.upsert(vectors=[{
            "id": f"{user_id}_{filename}_{i}",
            "values": embedding,
            "metadata": {
                "text": chunk,
                "user_id": user_id,
                "filename": filename,
                "chunk_id": i
            }
        }])

def process_file(contents, filename, user_id):
    try:
        logging.info(f"Processing file {filename} for user {user_id}")
        text = parse_file(contents, filename)
        if not text.strip():
            send_notification(user_id, f"Failed to extract text from '{filename}'.")
            return
        chunks = chunk_text(text)
        if not chunks:
            send_notification(user_id, f"No content could be extracted from '{filename}'.")
            return
        store_embeddings(user_id, filename, chunks)
        send_notification(user_id, f"Your file '{filename}' has been processed successfully with {len(chunks)} chunks.")
    except Exception as e:
        logging.error(f"Error processing file {filename}: {e}")
        send_notification(user_id, f"Error processing file '{filename}': {str(e)}")
        raise

from llama_index.embeddings.gemini import GeminiEmbedding

embed_model = GeminiEmbedding(api_key=GEMINI_API_KEY, model_name="models/embedding-001")
pinecone_index = get_pinecone_index()
